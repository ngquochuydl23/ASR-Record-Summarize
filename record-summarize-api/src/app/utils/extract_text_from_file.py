import os
import docx
import re
from PyPDF2 import PdfReader
from pptx import Presentation
import io


def normalize_text(text: str) -> str:
    text = re.sub(r'(?<!\n)\n(?!\n)', ' ', text)
    text = re.sub(r'\n{2,}', '\n', text)
    return text.strip()

def extract_text_from_pdf(path: str) -> str:
    reader = PdfReader(path)
    raw_text = " ".join([page.extract_text() for page in reader.pages if page.extract_text()])
    return normalize_text(raw_text)

def _extract_text_from_pdf_bytes(raw_bytes: bytes) -> str:
    reader = PdfReader(io.BytesIO(raw_bytes))
    raw_text = " ".join([page.extract_text() for page in reader.pages if page.extract_text()])
    return normalize_text(raw_text)

def extract_text_from_docx(path: str) -> str:
    doc = docx.Document(path)
    raw_text = " ".join([para.text for para in doc.paragraphs if para.text.strip()])
    return normalize_text(raw_text)

def _extract_text_from_docx_bytes(raw_bytes: bytes) -> str:
    doc = docx.Document(io.BytesIO(raw_bytes))
    raw_text = " ".join([para.text for para in doc.paragraphs if para.text.strip()])
    return normalize_text(raw_text)

def _extract_text_from_pptx_file(path: str) -> str:
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return normalize_text(" ".join(texts))

def _extract_text_from_pptx_bytes(raw_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(raw_bytes))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return normalize_text(" ".join(texts))

def extract_text(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(path)
    elif ext == ".docx":
        return extract_text_from_docx(path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

def extract_text_from_bytes(raw_bytes: bytes, ext: str) -> str:
    ext = ext.lower()
    if ext == ".pdf":
        return _extract_text_from_pdf_bytes(raw_bytes)
    elif ext == ".docx":
        return _extract_text_from_docx_bytes(raw_bytes)
    elif ext == ".pptx":
        return _extract_text_from_pptx_bytes(raw_bytes)
    elif ext == ".txt" or ext == ".vtt":
        return normalize_text(raw_bytes.decode("utf-8"))
    else:
        raise ValueError(f"Unsupported file type: {ext}")